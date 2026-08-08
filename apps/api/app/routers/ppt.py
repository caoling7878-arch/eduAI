from __future__ import annotations

import io
import json
import re

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy import select
from sqlalchemy.orm import Session

from ..audit_util import write_audit
from ..db import get_db
from ..models import PptJob, User
from ..rbac import require_staff
from ..schemas import PptIn, PptOut

router = APIRouter(prefix="/ppt", tags=["ppt"])


def _slides(title: str, outline: str) -> list[dict[str, str]]:
    lines = [ln.strip("- ").strip() for ln in outline.splitlines() if ln.strip()]
    if not lines:
        lines = ["背景与目标", "核心概念", "例题讲解", "课堂练习", "总结与作业"]
    slides = [{"title": title, "body": "课程导入与学习目标"}]
    for i, ln in enumerate(lines, 1):
        slides.append({"title": f"{i}. {ln}", "body": f"围绕「{ln}」展开讲解、举例与互动提问。"})
    slides.append({"title": "课堂小结", "body": "回顾要点，布置巩固练习。"})
    return slides


def _out(job: PptJob) -> PptOut:
    try:
        slides = json.loads(job.result_json or "[]")
    except json.JSONDecodeError:
        slides = []
    return PptOut(
        id=job.id,
        title=job.title,
        outline=job.outline,
        slides=slides if isinstance(slides, list) else [],
        status=job.status,
        created_at=job.created_at,
    )


def _build_pptx(title: str, slides: list[dict]) -> bytes:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise HTTPException(status_code=500, detail="服务未安装 python-pptx，无法导出") from e

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            prs.slide_width,
            Inches(0.18),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x0F, 0x6B, 0x5C)
        bar.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.2 if i == 0 else 1.2), Inches(12), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(s.get("title") or title)
        p.font.size = Pt(36 if i == 0 else 28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x2E, 0x28)
        p.alignment = PP_ALIGN.LEFT

        body = str(s.get("body") or "")
        if body:
            body_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(3.4 if i == 0 else 2.6), Inches(11.8), Inches(3.5)
            )
            btf = body_box.text_frame
            btf.word_wrap = True
            bp = btf.paragraphs[0]
            bp.text = body
            bp.font.size = Pt(18)
            bp.font.color.rgb = RGBColor(0x3A, 0x4A, 0x45)

        foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12), Inches(0.4))
        fp = foot.text_frame.paragraphs[0]
        fp.text = f"eduAI · {title} · {i + 1}/{len(slides)}"
        fp.font.size = Pt(12)
        fp.font.color.rgb = RGBColor(0x7A, 0x8A, 0x84)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _safe_filename(name: str) -> str:
    cleaned = re.sub(r'[\\/:*?"<>|]+', "_", name).strip() or "slides"
    return cleaned[:80]


@router.get("", response_model=list[PptOut])
def list_jobs(_: User = Depends(require_staff), db: Session = Depends(get_db)) -> list[PptOut]:
    return [_out(j) for j in db.scalars(select(PptJob).order_by(PptJob.id.desc()))]


@router.post("/generate", response_model=PptOut)
def generate(
    body: PptIn,
    admin: User = Depends(require_staff),
    db: Session = Depends(get_db),
) -> PptOut:
    slides = _slides(body.title, body.outline)
    job = PptJob(
        title=body.title,
        outline=body.outline,
        result_json=json.dumps(slides, ensure_ascii=False),
        status="done",
        created_by=admin.id,
    )
    db.add(job)
    write_audit(db, user=admin, action="ppt.generate", resource=body.title)
    db.commit()
    db.refresh(job)
    return _out(job)


@router.get("/{jid}", response_model=PptOut)
def get_job(jid: int, _: User = Depends(require_staff), db: Session = Depends(get_db)) -> PptOut:
    job = db.get(PptJob, jid)
    if not job:
        raise HTTPException(status_code=404, detail="任务不存在")
    return _out(job)


@router.get("/{jid}/export")
def export_pptx(
    jid: int,
    admin: User = Depends(require_staff),
    db: Session = Depends(get_db),
) -> StreamingResponse:
    job = db.get(PptJob, jid)
    if not job:
        raise HTTPException(status_code=404, detail="任务不存在")
    try:
        slides = json.loads(job.result_json or "[]")
    except json.JSONDecodeError:
        slides = []
    if not isinstance(slides, list) or not slides:
        raise HTTPException(status_code=400, detail="无可导出的幻灯片")
    data = _build_pptx(job.title, slides)
    write_audit(db, user=admin, action="ppt.export", resource=f"ppt:{jid}")
    db.commit()
    filename = f"{_safe_filename(job.title)}.pptx"
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
